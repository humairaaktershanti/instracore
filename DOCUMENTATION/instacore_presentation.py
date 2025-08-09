from pptx import Presentation

# Create a presentation
prs = Presentation()

slides_content = [
    {
        "title": "InstaCore: Institute Management System",
        "subtitle": "Presented by DiTech Squad",
        "bullet_points": []
    },
    {
        "title": "What is InstaCore?",
        "bullet_points": [
            "Web-based platform for institute management",
            "Modern, secure, and responsive",
            "Solves problems in outdated systems"
        ]
    },
    {
        "title": "Features – Phase 1",
        "bullet_points": [
            "Landing page with login & signup",
            "Custom Django authentication",
            "Email verification for signup",
            "OTP verification before account deletion",
            "Individual user dashboards",
            "Profile management (edit, password, delete)",
            "Role-aware navigation bar and footer"
        ]
    },
    {
        "title": "How It Works – User Flow",
        "bullet_points": [
            "Visit index page → Login or Signup",
            "Signup → Email verification → Login",
            "Dashboard → Profile management",
            "Delete account with OTP",
            "Secure logout"
        ]
    },
    {
        "title": "Tech Stack",
        "bullet_points": [
            "Frontend: HTML, CSS, JavaScript",
            "Backend: Django (Custom Auth, Email/OTP)",
            "Database: SQLite (dev) → PostgreSQL (prod)",
            "Planned: Django REST Framework for APIs"
        ]
    },
    {
        "title": "Project Timeline",
        "bullet_points": [
            "11 Aug – Present idea & mockups",
            "17 Aug – Frontend & auth completed",
            "21 Aug – CRUD operations & API integration",
            "24 Aug – Final submission with features"
        ]
    },
    {
        "title": "Future Plans",
        "bullet_points": [
            "Role-based dashboards (Admin, Teacher, Student)",
            "Institute-wide announcements",
            "Attendance & results module",
            "Public API documentation"
        ]
    },
    {
        "title": "UI Mockups",
        "bullet_points": [
            "Landing page",
            "Login page",
            "Dashboard",
            "Profile page"
        ]
    },
    {
        "title": "Final Thoughts",
        "bullet_points": [
            "Secure, flexible, and real-world ready",
            "Meets academic & practical needs",
            "Ready for future upgrades"
        ]
    },
    {
        "title": "Questions?",
        "bullet_points": [
            "We’re happy to answer your questions."
        ]
    }
]

# Function to add slides
for slide in slides_content:
    layout = prs.slide_layouts[1] if slide["bullet_points"] else prs.slide_layouts[0]
    sld = prs.slides.add_slide(layout)
    sld.shapes.title.text = slide["title"]
    if slide.get("subtitle"):
        sld.placeholders[1].text = slide["subtitle"]
    elif slide.get("bullet_points"):
        content = sld.placeholders[1].text_frame
        content.clear()
        for bp in slide["bullet_points"]:
            content.add_paragraph().text = bp

# Save file
prs.save("InstaCore_Presentation_DiTech_Squad.pptx")
print("✅ Presentation created successfully!")
